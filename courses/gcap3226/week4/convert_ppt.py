#!/usr/bin/env python3
"""
Convert PowerPoint to Markdown
Extracts text content from PPTX and formats it as markdown
"""

from pptx import Presentation
import sys
import os

def extract_text_from_slide(slide):
    """Extract all text from a slide"""
    text_runs = []
    
    # Get title
    if slide.shapes.title:
        title = slide.shapes.title.text.strip()
        if title:
            text_runs.append(('title', title))
    
    # Get text from all shapes
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            # Skip if it's the title we already got
            if slide.shapes.title and shape == slide.shapes.title:
                continue
            text_runs.append(('body', shape.text.strip()))
    
    return text_runs

def pptx_to_markdown(pptx_path, output_path=None):
    """Convert PPTX to Markdown"""
    
    # Load presentation
    prs = Presentation(pptx_path)
    
    # Prepare output
    if output_path is None:
        base_name = os.path.splitext(pptx_path)[0]
        output_path = f"{base_name}.md"
    
    markdown_content = []
    
    # Add header
    filename = os.path.basename(pptx_path)
    markdown_content.append(f"# {os.path.splitext(filename)[0]}")
    markdown_content.append("")
    markdown_content.append("---")
    markdown_content.append("")
    
    # Process each slide
    for slide_num, slide in enumerate(prs.slides, start=1):
        text_runs = extract_text_from_slide(slide)
        
        if not text_runs:
            continue
        
        markdown_content.append(f"## Slide {slide_num}")
        markdown_content.append("")
        
        for text_type, text in text_runs:
            if text_type == 'title':
                markdown_content.append(f"### {text}")
                markdown_content.append("")
            else:
                # Split by lines and handle bullet points
                lines = text.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        # Check if it looks like a bullet point
                        if line.startswith('•') or line.startswith('-') or line.startswith('*'):
                            markdown_content.append(f"- {line.lstrip('•-* ')}")
                        else:
                            markdown_content.append(line)
                markdown_content.append("")
        
        markdown_content.append("---")
        markdown_content.append("")
    
    # Write to file
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(markdown_content))
    
    print(f"✅ Converted {pptx_path}")
    print(f"📄 Output saved to: {output_path}")
    print(f"📊 Total slides processed: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    pptx_file = "Case Study II_ Evaluating the Efficiency of Bus Route Adjustments_ City bus No. 56.pptx"
    
    if os.path.exists(pptx_file):
        output_file = pptx_to_markdown(pptx_file)
    else:
        print(f"❌ File not found: {pptx_file}")
        sys.exit(1)
